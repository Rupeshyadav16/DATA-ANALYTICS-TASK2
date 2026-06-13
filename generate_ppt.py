import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

print("Creating automated PowerPoint BI Dashboard Mock-up...")

# Initialize Presentation and set to widescreen 16:9 layout
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use a blank slide layout
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# -----------------------------------------------------------------------------
# 1. HEADER SECTION (Title & Subtitle)
# -----------------------------------------------------------------------------
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Executive Business Intelligence Dashboard - Task 2"
p.font.size = Pt(28)
p.font.bold = True
p.font.name = "Arial"
p.font.color.rgb = RGBColor(26, 54, 93)  # Dark Corporate Navy Blue

p2 = tf.add_paragraph()
p2.text = "Data Analyst Intern: Rupesh Kumar Yadav  |  Company: ApexPlanet"
p2.font.size = Pt(13)
p2.font.italic = True
p2.font.name = "Arial"
p2.font.color.rgb = RGBColor(113, 128, 150) # Cool Muted Slate Gray

# -----------------------------------------------------------------------------
# 2. TOP KPI CARDS SECTION (3 Horizontal Summary Cards)
# -----------------------------------------------------------------------------
kpis = [
    {"title": "TOTAL SALES REVENUE", "value": "₹139.40 Million", "left": Inches(0.5)},
    {"title": "TOTAL SYSTEM VOLUME", "value": "1,000 Orders", "left": Inches(4.7)},
    {"title": "AVERAGE BASKET VALUE", "value": "₹139,399.44", "left": Inches(8.9)}
]

for kpi in kpis:
    # Draw standard background card block
    shape = slide.shapes.add_shape(
        1,  # Rectangular geometry ID
        kpi["left"], Inches(1.5), Inches(3.93), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(247, 250, 252) # Off-White Slate Tone
    shape.line.color.rgb = RGBColor(226, 232, 240)      # Clean subtle border
    
    # Text formatting inside cards
    ktf = shape.text_frame
    ktf.word_wrap = True
    
    kp1 = ktf.paragraphs[0]
    kp1.alignment = PP_ALIGN.CENTER
    kp1.text = kpi["title"]
    kp1.font.size = Pt(9)
    kp1.font.name = "Arial"
    kp1.font.color.rgb = RGBColor(160, 174, 192)
    
    kp2 = ktf.add_paragraph()
    kp2.alignment = PP_ALIGN.CENTER
    kp2.text = kpi["value"]
    kp2.font.size = Pt(18)
    kp2.font.bold = True
    kp2.font.name = "Arial"
    kp2.font.color.rgb = RGBColor(45, 55, 72)

# -----------------------------------------------------------------------------
# 3. GRAPHICS DATA CLUSTERING SECTION (Inserting Generated Python Charts)
# -----------------------------------------------------------------------------
chart1_path = "2_top_5_products_revenue.png"
chart2_path = "3_multivariate_correlation_heatmap.png"

# Left Chart: Top Products Breakdown
if os.path.exists(chart1_path):
    slide.shapes.add_picture(chart1_path, Inches(0.5), Inches(2.8), width=Inches(6.0), height=Inches(3.3))
    print("✓ Linked Top Products revenue graph layout.")
else:
    print(f"⚠ Warning: {chart1_path} not found. Running script without visual bindings.")

# Right Chart: Matrix Heatmap Correlation
if os.path.exists(chart2_path):
    slide.shapes.add_picture(chart2_path, Inches(6.833), Inches(2.8), width=Inches(6.0), height=Inches(3.3))
    print("✓ Linked Multi-feature Heatmap correlation graph layout.")
else:
    print(f"⚠ Warning: {chart2_path} not found. Running script without visual bindings.")

# -----------------------------------------------------------------------------
# 4. EXECUTIVE EXECUTIVE INSIGHT REPOSITORY (Bottom Banner)
# -----------------------------------------------------------------------------
footer_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.9))
footer_shape.fill.solid()
footer_shape.fill.fore_color.rgb = RGBColor(235, 248, 250) # Very Light Teal Hue
footer_shape.line.color.rgb = RGBColor(178, 231, 235)

ftf = footer_shape.text_frame
ftf.word_wrap = True

fp = ftf.paragraphs[0]
fp.text = "💡 Strategic Enterprise Metrics Conclusion Summary:"
fp.font.size = Pt(11)
fp.font.bold = True
fp.font.color.rgb = RGBColor(35, 78, 82)

fp2 = ftf.add_paragraph()
fp2.text = "• Core Revenue Drivers: Laptop (₹25.44M) and Mobile (₹25.33M) heavily out-scale alternative categories.\n• Data Pipeline Consistency: Features interaction displays uniform linear scaling matrices checking safely across unit price structures."
fp2.font.size = Pt(10.5)
fp2.font.color.rgb = RGBColor(45, 55, 72)

# Save the configured PowerPoint presentation layout
output_file = "Dashboard_Mockup.pptx"
prs.save(output_file)
print(f"\n🎉 Success! '{output_file}' has been automatically generated in your folder.")